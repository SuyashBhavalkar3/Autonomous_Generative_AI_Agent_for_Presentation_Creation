import logging
from pathlib import Path
import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import uuid

logger = logging.getLogger(__name__)


def download_image(url: str, dest_folder: Path, timeout: int = 10) -> Path | None:
	"""Download image to dest_folder and return file path, or None on failure."""
	if not url or not isinstance(url, str) or not url.startswith(("http://", "https://")):
		logger.warning("Invalid image URL, skipping download: %s", url)
		return None

	try:
		dest_folder.mkdir(parents=True, exist_ok=True)
		fname = url.split("/")[-1].split("?")[0]
		if not fname.lower().endswith((".jpg", ".jpeg", ".png")):
			fname = fname + ".jpg"
		out = dest_folder / fname
		with httpx.Client(timeout=timeout) as client:
			r = client.get(url)
			if r.status_code == 200 and r.content:
				out.write_bytes(r.content)
				return out
			else:
				logger.warning("Image download failed %s status=%s", url, r.status_code)
				return None
	except Exception as e:
		logger.warning("Image download exception %s: %s", url, e)
		return None


def _apply_gamma_theme(prs: Presentation):
	"""Apply a simple dark Gamma-like theme to the Presentation.

	- Dark background
	- Bold white titles
	- Muted gray bullets
	"""
	# Attempt to set slide background for existing slides (new slides will get style from here)
	for slide in prs.slides:
		try:
			fill = slide.background.fill
			fill.solid()
			fill.fore_color.rgb = RGBColor(18, 18, 18)  # very dark gray
		except Exception:
			continue


def build_presentation(slides: list[dict], out_path: Path | str) -> Path:
    """Create PPTX with title, bullets, and images below text, fully working."""
    prs = Presentation()
    tmp_dir = Path("output") / "images" / f"ppt_builder_{uuid.uuid4().hex}"
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    left_margin = Inches(0.5)
    right_margin = Inches(0.5)
    bottom_margin = Inches(0.5)
    top_margin = Inches(0.3)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for s in slides:
        title = s.get("title", "Untitled")
        bullets = s.get("bullets", [])[:4]

        # --- Image fetch (working logic from commented version) ---
        image_path = s.get("image_path")
        if not image_path:
            image_url = s.get("image_url")
            if image_url:
                try:
                    _downloaded = download_image(image_url, tmp_dir)
                    image_path = str(_downloaded) if _downloaded else None
                except Exception:
                    image_path = None

        # --- Add slide ---
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(18, 18, 18)
        except Exception:
            pass

        # --- Title ---
        title_height = Inches(1.0)
        try:
            title_box = slide.shapes.add_textbox(left_margin, top_margin,
													slide_w - left_margin - right_margin,
													title_height)
            tf = title_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(255, 255, 255)
        except Exception:
            logger.debug("Failed to add title box")

        # --- Bullets ---
        body_top = top_margin + title_height + Inches(0.2)
        body_height = Inches(3.0)
        try:
            body_box = slide.shapes.add_textbox(left_margin, body_top,
												slide_w - left_margin - right_margin,
												body_height)
            tf = body_box.text_frame
            tf.clear()
            for idx, b in enumerate(bullets):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(200, 200, 200)
        except Exception:
            logger.debug("Failed to add bullets")

        # --- Image below text (fully working from commented version) ---
        if image_path:
            try:
                img_file = Path(image_path)
                if img_file.exists() and img_file.stat().st_size > 0:
                    # Place image just below the text box (body_top + body_height + padding)
                    pic_left = left_margin
                    pic_top = body_top + body_height + Inches(0.2)
                    # Set width to slide width minus margins; aspect ratio preserved
                    pic_width = slide_w - left_margin - right_margin
                    slide.shapes.add_picture(str(img_file), pic_left, pic_top, width=pic_width)
                else:
                    logger.warning("Image file missing or empty: %s", image_path)
            except Exception as e:
                logger.warning("Embedding image failed: %s", e)

    # Save PPT
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out

class PPTBuilder:
	"""Small wrapper class for compatibility with older tests/code.

	Usage:
	    builder = PPTBuilder(output_path="output/foo.pptx")
	    builder.build(slide_agent_output)
	"""

	def __init__(self, output_path: str | Path):
		self.output_path = Path(output_path)

	def build(self, slide_agent_output: dict | list) -> Path:
		# Accept either a dict with `slides` key or a direct list
		if isinstance(slide_agent_output, dict):
			slides = slide_agent_output.get("slides") or []
		else:
			slides = slide_agent_output
		# Defensive: ensure slides is a list of dict
		if not isinstance(slides, list):
			raise ValueError("Invalid slides input for PPTBuilder")
		return build_presentation(slides, self.output_path)